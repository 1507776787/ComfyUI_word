import os
import re
import json
from datetime import datetime
from urllib.parse import quote
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Pt as PptPt, Inches as PptInches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from html.parser import HTMLParser
import folder_paths

try:
    # Compatible path for most python-pptx versions
    from pptx.oxml.xmlchemy import OxmlElement
except Exception:
    OxmlElement = None


OUTPUT_SUBFOLDER = "ComfyUI_word"


def _normalize_base_url(base_url):
    base_url = (base_url or "").strip()
    if not base_url:
        return ""
    return base_url.rstrip("/")


def _get_unique_output_path(file_name, extension):
    safe_name = os.path.basename((file_name or "output").strip()) or "output"
    if not safe_name.lower().endswith(extension):
        safe_name = safe_name + extension

    output_dir = os.path.join(folder_paths.get_output_directory(), OUTPUT_SUBFOLDER)
    os.makedirs(output_dir, exist_ok=True)

    file_path = os.path.join(output_dir, safe_name)
    counter = 1
    while os.path.exists(file_path):
        name, ext = os.path.splitext(safe_name)
        unique_name = f"{name}_{counter}{ext}"
        file_path = os.path.join(output_dir, unique_name)
        counter += 1
    return file_path


def _build_download_url(file_path, base_url=""):
    filename = os.path.basename(file_path)
    query = f"filename={quote(filename)}&type=output&subfolder={quote(OUTPUT_SUBFOLDER)}"
    relative_url = f"/view?{query}"
    normalized = _normalize_base_url(base_url)
    if normalized:
        return f"{normalized}{relative_url}"
    return relative_url


class TextToWordNode:
    
    def __init__(self):
        self.output_dir = os.path.join(folder_paths.get_output_directory(), OUTPUT_SUBFOLDER)
        os.makedirs(self.output_dir, exist_ok=True)
    
    @classmethod
    def INPUT_TYPES(cls):
        return {
            "required": {
                "文字内容": ("STRING", {
                    "multiline": True,
                    "default": "请输入文字内容..."
                }),
                "文件名": ("STRING", {
                    "default": "output"
                }),
            },
            "optional": {
                "字体名称": ("STRING", {
                    "default": "微软雅黑"
                }),
                "字体大小": ("INT", {
                    "default": 12,
                    "min": 8,
                    "max": 72
                }),
                "对齐方式": (["左对齐", "居中", "右对齐"], {
                    "default": "左对齐"
                }),
                "行间距": ("FLOAT", {
                    "default": 1.5,
                    "min": 1.0,
                    "max": 3.0,
                    "step": 0.1
                }),
                "下载地址前缀": ("STRING", {
                    "default": ""
                }),
            }
        }
    
    RETURN_TYPES = ("STRING", "STRING")
    RETURN_NAMES = ("文件路径", "下载链接")
    OUTPUT_NODE = True
    FUNCTION = "generate_word"
    CATEGORY = "liujian"
    
    def generate_word(self, 文字内容, 文件名, 字体名称="微软雅黑", 字体大小=12,
                     对齐方式="左对齐", 行间距=1.5, 下载地址前缀=""):
        doc = Document()
        
        paragraph = doc.add_paragraph()
        
        alignment_map = {
            "左对齐": WD_ALIGN_PARAGRAPH.LEFT,
            "居中": WD_ALIGN_PARAGRAPH.CENTER,
            "右对齐": WD_ALIGN_PARAGRAPH.RIGHT
        }
        paragraph.alignment = alignment_map.get(对齐方式, WD_ALIGN_PARAGRAPH.LEFT)
        
        paragraph_format = paragraph.paragraph_format
        paragraph_format.line_spacing = 行间距
        
        run = paragraph.add_run(文字内容)
        run.font.name = 字体名称
        run.font.size = Pt(字体大小)
        
        run._element.rPr.rFonts.set(
            '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia',
            字体名称
        )
        
        file_path = _get_unique_output_path(文件名, ".docx")
        
        doc.save(file_path)
        download_url = _build_download_url(file_path, 下载地址前缀)
        return {
            "ui": {"text": [f"本地文件: {file_path}", f"下载链接: {download_url}"]},
            "result": (file_path, download_url)
        }
    
class TextToWordAdvancedNode:
    
    def __init__(self):
        self.output_dir = os.path.join(folder_paths.get_output_directory(), OUTPUT_SUBFOLDER)
        os.makedirs(self.output_dir, exist_ok=True)
    
    @classmethod
    def INPUT_TYPES(cls):
        return {
            "required": {
                "文字内容": ("STRING", {
                    "multiline": True,
                    "default": "请输入文字内容...\n使用空行分隔段落"
                }),
                "文件名": ("STRING", {
                    "default": "output"
                }),
            },
            "optional": {
                "标题": ("STRING", {
                    "default": ""
                }),
                "标题字号": ("INT", {
                    "default": 18,
                    "min": 12,
                    "max": 48
                }),
                "正文字体": ("STRING", {
                    "default": "微软雅黑"
                }),
                "正文字号": ("INT", {
                    "default": 12,
                    "min": 8,
                    "max": 72
                }),
                "对齐方式": (["左对齐", "居中", "右对齐", "两端对齐"], {
                    "default": "左对齐"
                }),
                "行间距": ("FLOAT", {
                    "default": 1.5,
                    "min": 1.0,
                    "max": 3.0,
                    "step": 0.1
                }),
                "下载地址前缀": ("STRING", {
                    "default": ""
                }),
            }
        }
    
    RETURN_TYPES = ("STRING", "STRING")
    RETURN_NAMES = ("文件路径", "下载链接")
    OUTPUT_NODE = True
    FUNCTION = "generate_word_advanced"
    CATEGORY = "liujian"
    
    def generate_word_advanced(self, 文字内容, 文件名, 标题="", 标题字号=18,
                                正文字体="微软雅黑", 正文字号=12,
                                对齐方式="左对齐", 行间距=1.5, 下载地址前缀=""):
        doc = Document()
        
        alignment_map = {
            "左对齐": WD_ALIGN_PARAGRAPH.LEFT,
            "居中": WD_ALIGN_PARAGRAPH.CENTER,
            "右对齐": WD_ALIGN_PARAGRAPH.RIGHT,
            "两端对齐": WD_ALIGN_PARAGRAPH.JUSTIFY
        }
        
        if 标题:
            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run(标题)
            title_run.font.name = 正文字体
            title_run.font.size = Pt(标题字号)
            title_run.bold = True
            title_run._element.rPr.rFonts.set(
                '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia',
                正文字体
            )
            doc.add_paragraph()
        
        paragraphs_text = [p.strip() for p in 文字内容.split('\n\n') if p.strip()]
        
        for para_text in paragraphs_text:
            paragraph = doc.add_paragraph()
            paragraph.alignment = alignment_map.get(对齐方式, WD_ALIGN_PARAGRAPH.LEFT)
            
            paragraph_format = paragraph.paragraph_format
            paragraph_format.line_spacing = 行间距
            paragraph_format.first_line_indent = Inches(0.3)
            
            run = paragraph.add_run(para_text)
            run.font.name = 正文字体
            run.font.size = Pt(正文字号)
            run._element.rPr.rFonts.set(
                '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia',
                正文字体
            )
        
        file_path = _get_unique_output_path(文件名, ".docx")
        
        doc.save(file_path)
        download_url = _build_download_url(file_path, 下载地址前缀)
        return {
            "ui": {"text": [f"本地文件: {file_path}", f"下载链接: {download_url}"]},
            "result": (file_path, download_url)
        }
    
class SlideHTMLParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.slides = []
        self.current_slide = None
        self.current_element = None
        self.current_text = ""
        self.current_attrs = {}
        self.list_items = []
        self.in_list = False
        self.current_region = None
        self.current_section_stack = []
        self.section_counter = 0
        self.global_theme = None
        self.global_footer = None
    
    def _current_section(self):
        return self.current_section_stack[-1] if self.current_section_stack else None
    
    def _resolve_region(self, attrs_dict):
        if attrs_dict.get("region"):
            return attrs_dict.get("region")
        section = self._current_section()
        if section and section.get("attrs", {}).get("region"):
            return section["attrs"]["region"]
        return self.current_region
    
    def _append_component(self, comp_type, attrs_dict):
        if not self.current_slide:
            return
        self.current_slide["content"].append({
            "type": comp_type,
            "attrs": attrs_dict,
            "region": self._resolve_region(attrs_dict),
            "section": self._current_section()
        })
        
    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        
        if tag == "theme":
            if self.current_slide:
                self.current_slide["theme"] = attrs_dict.get("name", attrs_dict.get("preset", ""))
            else:
                self.global_theme = attrs_dict.get("name", attrs_dict.get("preset", ""))
            return
        
        if tag == "slide":
            self.current_slide = {
                "title": None,
                "title_style": {},
                "content": [],
                "background": None,
                "layout": "default",
                "theme": self.global_theme,
                "footer": self.global_footer
            }
            self.current_region = None
            self.current_section_stack = []
        elif tag == "layout":
            if self.current_slide:
                self.current_slide["layout"] = attrs_dict.get("name", "default")
        elif tag == "left":
            self.current_region = "left"
        elif tag == "right":
            self.current_region = "right"
        elif tag == "section":
            if self.current_slide:
                self.section_counter += 1
                section_obj = {"id": self.section_counter, "attrs": attrs_dict}
                self.current_section_stack.append(section_obj)
        elif tag == "title":
            self.current_element = "title"
            self.current_attrs = attrs_dict
            self.current_text = ""
        elif tag == "h1":
            self.current_element = "h1"
            self.current_attrs = attrs_dict
            self.current_text = ""
        elif tag == "h2":
            self.current_element = "h2"
            self.current_attrs = attrs_dict
            self.current_text = ""
        elif tag == "p":
            self.current_element = "p"
            self.current_attrs = attrs_dict
            self.current_text = ""
        elif tag == "ul":
            self.in_list = True
            self.list_items = []
        elif tag == "ol":
            self.in_list = True
            self.list_items = []
            self.list_ordered = True
        elif tag == "li":
            self.current_element = "li"
            self.current_attrs = attrs_dict
            self.current_text = ""
        elif tag == "card":
            self._append_component("card", attrs_dict)
        elif tag == "badge":
            self._append_component("badge", attrs_dict)
        elif tag == "button":
            self._append_component("button", attrs_dict)
        elif tag == "icon":
            self._append_component("icon", attrs_dict)
        elif tag == "chart":
            self._append_component("chart", attrs_dict)
        elif tag == "table":
            self._append_component("table", attrs_dict)
        elif tag == "callout":
            self._append_component("callout", attrs_dict)
        elif tag == "divider":
            self._append_component("divider", attrs_dict)
        elif tag == "spacer":
            self._append_component("spacer", attrs_dict)
        elif tag == "toc":
            self._append_component("toc", attrs_dict)
        elif tag == "footer":
            if self.current_slide:
                self.current_slide["footer"] = attrs_dict
            else:
                self.global_footer = attrs_dict
        elif tag == "background":
            if self.current_slide:
                self.current_slide["background"] = attrs_dict.get("color", "#FFFFFF")
    
    def handle_endtag(self, tag):
        if tag == "slide" and self.current_slide:
            self.slides.append(self.current_slide)
            self.current_slide = None
            self.current_region = None
            self.current_section_stack = []
        elif tag in ["left", "right"]:
            self.current_region = None
        elif tag == "section":
            if self.current_section_stack:
                self.current_section_stack.pop()
        elif tag == "title" and self.current_slide:
            self.current_slide["title"] = self.current_text
            self.current_slide["title_style"] = self.current_attrs
            self.current_element = None
        elif tag in ["h1", "h2"] and self.current_slide:
            content_item = {
                "type": tag,
                "text": self.current_text,
                "style": self.current_attrs,
                "region": self._resolve_region(self.current_attrs),
                "section": self._current_section()
            }
            self.current_slide["content"].append(content_item)
            self.current_element = None
        elif tag == "p" and self.current_slide:
            content_item = {
                "type": "p",
                "text": self.current_text,
                "style": self.current_attrs,
                "region": self._resolve_region(self.current_attrs),
                "section": self._current_section()
            }
            self.current_slide["content"].append(content_item)
            self.current_element = None
        elif tag == "li":
            self.list_items.append({
                "text": self.current_text,
                "style": self.current_attrs,
                "region": self._resolve_region(self.current_attrs),
                "section": self._current_section()
            })
            self.current_element = None
        elif tag in ["ul", "ol"] and self.current_slide:
            self.current_slide["content"].append({
                "type": "list",
                "items": self.list_items,
                "ordered": tag == "ol",
                "region": self._resolve_region({}),
                "section": self._current_section(),
                "attrs": {}
            })
            self.in_list = False
            self.list_items = []
    
    def handle_data(self, data):
        if self.current_element:
            self.current_text += data


class HTMLToPPTNode:
    
    def __init__(self):
        self.output_dir = os.path.join(folder_paths.get_output_directory(), OUTPUT_SUBFOLDER)
        os.makedirs(self.output_dir, exist_ok=True)
        self.themes_file = os.path.join(os.path.dirname(__file__), "themes.json")
        self.theme_presets = self._merge_theme_presets()
    
    @classmethod
    def INPUT_TYPES(cls):
        return {
            "required": {
                "HTML代码": ("STRING", {
                    "multiline": True,
                    "default": """<slide>
  <theme name="enterprise_blue" />
  <layout name="kpi-dashboard" />
  <background color="#F8FAFC" />
  <footer text="AI 报告 | 第 {page}/{total} 页 | {date}" />
  <title style="font-size:38;color:#1E3A8A">智能体工作流看板</title>
  <section region="left" gap="0.15">
    <badge text="LIVE" bg="#DBEAFE" color="#1D4ED8" />
    <h2 style="font-size:24;color:#0F172A">项目进度</h2>
    <divider color="#CBD5E1" />
    <ul>
      <li style="font-size:16;color:#334155">需求分析：100%</li>
      <li style="font-size:16;color:#334155">节点开发：80%</li>
      <li style="font-size:16;color:#334155">联调测试：60%</li>
    </ul>
    <spacer height="0.2" />
    <icon name="rocket" text="阶段目标推进中" color="#4F46E5" />
  </section>
  <section region="right" gap="0.2">
    <card bg="#E0E7FF" border="#4F46E5" title="本周 KPI" body="生成成功率 98.6%&#10;平均耗时 6.2s" />
    <callout level="success" title="稳定性" body="任务成功率连续 7 天 > 98%" />
    <button text="查看下载链接" bg="#4F46E5" color="#FFFFFF" />
  </section>
</slide>
<slide>
  <theme name="modern_dark" />
  <layout name="comparison" />
  <footer text="方案对比 | 第 {page}/{total} 页" />
  <background color="#0F172A" />
  <title style="font-size:34;color:#E2E8F0">能力矩阵</title>
  <section region="left" gap="0.15">
    <callout level="info" title="现状" body="人工排版耗时高，交付不稳定" />
    <chart type="bar" title="效率提升(%)" data="文档:62,演示:55,下载链路:88" bar-color="#22D3EE" />
  </section>
  <section region="right" gap="0.15">
    <table columns="维度,改造前,改造后" rows="生成耗时|15min|4min;下载方式|本地路径|直链下载;稳定性|人工流程|自动化流程" />
    <button text="批量生成" bg="#22C55E" />
  </section>
</slide>
<slide>
  <theme name="enterprise_blue" />
  <layout name="default" />
  <title style="font-size:34">目录</title>
  <toc />
</slide>
<slide>
  <theme name="soft_gray" />
  <background color="#FFFFFF" />
  <title style="font-size:36;color:#1E3A8A">下载能力演示</title>
  <p style="font-size:18;color:#374151">文件落盘到 output/ComfyUI_word，并直接返回下载链接。</p>
  <section gap="0.12">
    <badge text="NEW" bg="#DBEAFE" color="#1D4ED8" />
    <button text="导出报告" bg="#0EA5E9" color="#FFFFFF" />
    <button text="批量生成" bg="#22C55E" color="#FFFFFF" />
    <divider />
    <p style="font-size:16;color:#475569">本页全部元素仅用 HTML 生成，无需图片和视频文件。</p>
  </section>
</slide>"""
                }),
                "文件名": ("STRING", {
                    "default": "output"
                }),
            },
            "optional": {
                "默认字体": ("STRING", {
                    "default": "微软雅黑"
                }),
                "幻灯片宽度": ("FLOAT", {
                    "default": 12.8,
                    "min": 8.0,
                    "max": 30.0,
                    "step": 0.1
                }),
                "幻灯片高度": ("FLOAT", {
                    "default": 7.2,
                    "min": 4.5,
                    "max": 20.0,
                    "step": 0.1
                }),
                "下载地址前缀": ("STRING", {
                    "default": ""
                }),
            }
        }
    
    RETURN_TYPES = ("STRING", "STRING")
    RETURN_NAMES = ("文件路径", "下载链接")
    OUTPUT_NODE = True
    FUNCTION = "generate_ppt"
    CATEGORY = "liujian"

    BUILTIN_THEME_PRESETS = {
        "enterprise_blue": {
            "background": "#F8FAFC",
            "title": "#1E3A8A",
            "text": "#334155",
            "muted": "#64748B",
            "card_bg": "#EEF2FF",
            "card_border": "#C7D2FE",
            "primary": "#4F46E5",
            "accent": "#0EA5E9",
            "margin_top": 0.5,
            "margin_bottom": 0.45,
            "line_gap": 0.14
        },
        "modern_dark": {
            "background": "#0F172A",
            "title": "#E2E8F0",
            "text": "#F1F5F9",
            "muted": "#CBD5E1",
            "card_bg": "#1E293B",
            "card_border": "#334155",
            "primary": "#22D3EE",
            "accent": "#A78BFA",
            "margin_top": 0.5,
            "margin_bottom": 0.45,
            "line_gap": 0.14
        },
        "soft_gray": {
            "background": "#FFFFFF",
            "title": "#111827",
            "text": "#374151",
            "muted": "#6B7280",
            "card_bg": "#F3F4F6",
            "card_border": "#D1D5DB",
            "primary": "#2563EB",
            "accent": "#10B981",
            "margin_top": 0.5,
            "margin_bottom": 0.45,
            "line_gap": 0.14
        },
        "modern_orange": {
            "background": "#FFF7ED",
            "title": "#7C2D12",
            "text": "#431407",
            "muted": "#9A3412",
            "card_bg": "#FFFFFF",
            "card_border": "#FED7AA",
            "primary": "#F97316",
            "accent": "#FB923C",
            "margin_top": 0.5,
            "margin_bottom": 0.45,
            "line_gap": 0.14
        },
        "apple_light": {
            "background": "#F5F6F8",
            "title": "#111827",
            "text": "#4B5563",
            "muted": "#6B7280",
            "card_bg": "#FFFFFF",
            "card_border": "#E5E7EB",
            "primary": "#111827",
            "accent": "#3B82F6",
            "margin_top": 0.5,
            "margin_bottom": 0.45,
            "line_gap": 0.14
        },
        "glass_cyan": {
            "background": "#0B1220",
            "title": "#F8FAFC",
            "text": "#C7D2FE",
            "muted": "#94A3B8",
            "card_bg": "#1E293BCC",
            "card_border": "#334155",
            "primary": "#22D3EE",
            "accent": "#38BDF8",
            "margin_top": 0.5,
            "margin_bottom": 0.45,
            "line_gap": 0.14
        },
        "violet_future": {
            "background": "#1E1B4B",
            "title": "#EDE9FE",
            "text": "#C4B5FD",
            "muted": "#A78BFA",
            "card_bg": "#312E81",
            "card_border": "#6366F1",
            "primary": "#8B5CF6",
            "accent": "#A78BFA",
            "margin_top": 0.5,
            "margin_bottom": 0.45,
            "line_gap": 0.14
        },
        "emerald_growth": {
            "background": "#ECFDF5",
            "title": "#065F46",
            "text": "#064E3B",
            "muted": "#047857",
            "card_bg": "#FFFFFF",
            "card_border": "#A7F3D0",
            "primary": "#10B981",
            "accent": "#34D399",
            "margin_top": 0.5,
            "margin_bottom": 0.45,
            "line_gap": 0.14
        },
        "rose_brand": {
            "background": "#FFF1F2",
            "title": "#9F1239",
            "text": "#881337",
            "muted": "#BE123C",
            "card_bg": "#FFFFFF",
            "card_border": "#FECDD3",
            "primary": "#E11D48",
            "accent": "#FB7185",
            "margin_top": 0.5,
            "margin_bottom": 0.45,
            "line_gap": 0.14
        },
        "midnight_cyan": {
            "background": "#020617",
            "title": "#E2E8F0",
            "text": "#BAE6FD",
            "muted": "#7DD3FC",
            "card_bg": "#0F172A",
            "card_border": "#1E293B",
            "primary": "#06B6D4",
            "accent": "#22D3EE",
            "margin_top": 0.5,
            "margin_bottom": 0.45,
            "line_gap": 0.14
        },
        "warm_business": {
            "background": "#FFFBEB",
            "title": "#78350F",
            "text": "#92400E",
            "muted": "#B45309",
            "card_bg": "#FFFFFF",
            "card_border": "#FCD34D",
            "primary": "#D97706",
            "accent": "#F59E0B",
            "margin_top": 0.5,
            "margin_bottom": 0.45,
            "line_gap": 0.14
        }
    }

    ICON_MAP = {
        "check": "✔",
        "warn": "⚠",
        "rocket": "🚀",
        "chart": "📊",
        "doc": "📄",
        "link": "🔗",
        "idea": "💡"
    }
    
    def parse_style(self, style_str):
        style_dict = {}
        if not style_str:
            return style_dict
        for item in style_str.split(';'):
            if ':' in item:
                key, value = item.split(':', 1)
                style_dict[key.strip()] = value.strip()
        return style_dict

    def _style_or_attr(self, attrs, style, key, default=None):
        if key in attrs:
            return attrs.get(key)
        return style.get(key, default)

    def _component_color(self, attrs, style, keys, default_color):
        for key in keys:
            # Support both kebab-case and snake_case from attrs/style.
            v = self._style_or_attr(attrs, style, key, None)
            if v is None:
                v = self._style_or_attr(attrs, style, key.replace("-", "_"), None)
            color = self._resolve_color(v, None)
            if color:
                return color
        return default_color

    def _parse_length(self, value, default, total=None):
        if value is None or str(value).strip() == "":
            return default
        s = str(value).strip().lower()
        try:
            if s.endswith('%') and total is not None:
                return total * float(s[:-1]) / 100.0
            if s.endswith('cm'):
                return float(s[:-2]) / 2.54
            if s.endswith('pt'):
                return float(s[:-2]) / 72.0
            if s.endswith('px'):
                return float(s[:-2]) / 96.0
            if s.endswith('in'):
                return float(s[:-2])
            return float(s)
        except ValueError:
            return default

    def _region_box(self, region, slide_width):
        margin = 0.6
        gap = 0.3
        col_width = (slide_width - margin * 2 - gap) / 2.0
        if region == "left":
            return margin, col_width
        if region == "right":
            return margin + col_width + gap, col_width
        return 0.5, slide_width - 1.0

    def _text_align(self, value, default=PP_ALIGN.LEFT):
        align = (value or "").strip().lower()
        if align in ("center", "middle", "居中"):
            return PP_ALIGN.CENTER
        if align in ("right", "end", "右对齐", "右"):
            return PP_ALIGN.RIGHT
        if align in ("justify", "两端对齐"):
            return PP_ALIGN.JUSTIFY
        if align in ("left", "start", "左对齐", "左"):
            return PP_ALIGN.LEFT
        return default

    def _to_bool(self, value, default=False):
        if value is None:
            return default
        return str(value).strip().lower() in {"1", "true", "yes", "y", "on"}

    def _normalize_theme(self, theme, fallback):
        norm = dict(fallback)
        if isinstance(theme, dict):
            norm.update(theme)
        return norm

    def _load_external_theme_presets(self):
        if not os.path.isfile(self.themes_file):
            return {}
        try:
            with open(self.themes_file, "r", encoding="utf-8") as f:
                data = json.load(f)
            if not isinstance(data, dict):
                return {}
            return {str(k): v for k, v in data.items() if isinstance(v, dict)}
        except Exception:
            return {}

    def _merge_theme_presets(self):
        base = dict(self.BUILTIN_THEME_PRESETS)
        external = self._load_external_theme_presets()
        fallback = base.get("enterprise_blue", next(iter(base.values())))

        for name, theme in external.items():
            if name in base:
                base[name] = self._normalize_theme(theme, base[name])
            else:
                base[name] = self._normalize_theme(theme, fallback)
        return base

    def _theme(self, name):
        key = (name or "enterprise_blue").strip()
        if key in self.theme_presets:
            return self.theme_presets[key]
        return self.theme_presets.get("enterprise_blue", next(iter(self.theme_presets.values())))

    def _draw_rect(self, slide, x, y, w, h, bg=None, border=None, rounded=True, shadow=False, radius=None):
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type, PptInches(x), PptInches(y), PptInches(w), PptInches(h)
        )
        if rounded and radius is not None:
            try:
                shape.adjustments[0] = float(radius)
            except Exception:
                pass

        if bg:
            shape.fill.solid()
            color_tuple = self.hex_to_rgb(bg)
            r, g, b = color_tuple[:3]
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
        else:
            shape.fill.background()
            
        if border:
            color_tuple = self.hex_to_rgb(border)
            r, g, b = color_tuple[:3]
            shape.line.color.rgb = RGBColor(r, g, b)
        else:
            shape.line.fill.background()
            
        return shape

    def _add_textbox(self, slide, x, y, w, h, text, font_name, size=14, color="#111827", bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
        tf = box.text_frame
        tf.word_wrap = True
        # Auto-fit text to reduce overflow when content is slightly long.
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        para = tf.paragraphs[0]
        para.text = text
        para.font.name = font_name
        para.font.size = PptPt(size)
        para.font.bold = bold
        para.alignment = align
        if color:
            color_tuple = self.hex_to_rgb(color)
            r, g, b = color_tuple[:3]
            para.font.color.rgb = RGBColor(r, g, b)
            if OxmlElement is not None and len(color_tuple) > 3 and color_tuple[3] < 255:
                opacity = int((color_tuple[3] / 255.0) * 100000)
                srgbClr = para.font.color._xClr
                alpha_elem = OxmlElement('a:alpha')
                alpha_elem.set('val', str(opacity))
                srgbClr.append(alpha_elem)
        return box

    def _parse_pairs(self, data_str):
        points = []
        for item in (data_str or "").split(","):
            item = item.strip()
            if not item:
                continue
            if ":" in item:
                k, v = item.split(":", 1)
                try:
                    points.append((k.strip(), float(v.strip())))
                except ValueError:
                    continue
        return points

    def _parse_rows(self, rows_str):
        rows = []
        for row in (rows_str or "").split(";"):
            row = row.strip()
            if not row:
                continue
            rows.append([c.strip() for c in row.split("|")])
        return rows

    def _render_footer(self, slide, footer_attrs, page_num, total_pages, slide_width, slide_height, theme, font_name, title_text):
        if not footer_attrs:
            return
        show = self._to_bool(footer_attrs.get("show", "true"), True)
        if not show:
            return
        raw = footer_attrs.get("text", "第 {page}/{total} 页")
        footer_text = raw.format(
            page=page_num,
            total=total_pages,
            date=datetime.now().strftime("%Y-%m-%d"),
            title=title_text or ""
        )
        x = self._parse_length(footer_attrs.get("x"), 0.5, slide_width)
        y = self._parse_length(footer_attrs.get("y"), slide_height - 0.42, slide_height)
        w = self._parse_length(footer_attrs.get("width"), slide_width - 1.0, slide_width)
        h = self._parse_length(footer_attrs.get("height"), 0.3, slide_height)
        color = self._resolve_color(footer_attrs.get("color"), theme["muted"])
        align = self._text_align(footer_attrs.get("text-align"), PP_ALIGN.RIGHT)
        self._add_textbox(slide, x, y, w, h, footer_text, font_name, size=11, color=color, align=align)

    def _ensure_section_state(self, section_state, section, slide, lane_y, lane, slide_width, slide_height, theme):
        if not section:
            return None
        sid = section["id"]
        if sid in section_state:
            return section_state[sid]

        attrs = section.get("attrs", {})
        sec_lane = attrs.get("region", lane) if attrs.get("region") in ("left", "right", "main") else lane
        lane_x, lane_w = self._region_box(sec_lane, slide_width)
        sec_x = self._parse_length(attrs.get("x"), lane_x, slide_width)
        sec_y = self._parse_length(attrs.get("y"), lane_y[sec_lane], slide_height)
        sec_w = self._parse_length(attrs.get("width"), lane_w, slide_width)
        sec_h = self._parse_length(attrs.get("height"), max(1.0, slide_height - sec_y - theme["margin_bottom"]), slide_height)
        padding = self._parse_length(attrs.get("padding"), 0.18)
        gap = self._parse_length(attrs.get("gap"), theme["line_gap"])

        sec_bg = self._resolve_color(attrs.get("bg"), None)
        sec_border = self._resolve_color(attrs.get("border"), None)
        shadow = self._to_bool(attrs.get("shadow"), False)
        rounded = self._to_bool(attrs.get("rounded"), False)
        radius = self._parse_length(attrs.get("radius"), None)
        if sec_bg or sec_border:
            self._draw_rect(slide, sec_x, sec_y, sec_w, sec_h, bg=sec_bg, border=sec_border, rounded=rounded, shadow=shadow, radius=radius)

        state = {
            "id": sid,
            "lane": sec_lane,
            "x": sec_x,
            "y": sec_y,
            "w": sec_w,
            "h": sec_h,
            "padding": padding,
            "gap": gap,
            "cursor": sec_y + padding
        }
        section_state[sid] = state
        lane_y[sec_lane] = max(lane_y[sec_lane], sec_y + padding)
        return state

    def _resolve_color(self, color_str, default_color=None):
        if not color_str:
            return default_color
        color_str = color_str.strip()
        if color_str.startswith('#'):
            hex_part = color_str[1:]
            if len(hex_part) == 3:
                hex_part = ''.join(ch * 2 for ch in hex_part)
            if len(hex_part) == 6 and re.fullmatch(r'[0-9a-fA-F]{6}', hex_part):
                return f"#{hex_part}"
            if len(hex_part) == 8 and re.fullmatch(r'[0-9a-fA-F]{8}', hex_part):
                return f"#{hex_part}"
            return default_color
        if color_str.startswith('rgba'):
            match = re.search(r'rgba\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*([\d.]+)\s*\)', color_str)
            if match:
                r, g, b, a = match.groups()
                a_int = int(float(a) * 255)
                return f"#{int(r):02x}{int(g):02x}{int(b):02x}{a_int:02x}"
        if color_str.startswith('rgb'):
            match = re.search(r'rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color_str)
            if match:
                r, g, b = match.groups()
                return f"#{int(r):02x}{int(g):02x}{int(b):02x}"
        return default_color
    
    def get_font_size(self, style, default_size=18):
        if 'font-size' in style:
            size_str = style['font-size']
            match = re.search(r'(\d+)', size_str)
            if match:
                return int(match.group(1))
        return default_size
    
    def get_color(self, style, default_color=None):
        if 'color' in style:
            return self._resolve_color(style['color'], default_color)
        return default_color
    
    def hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 8:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4, 6))
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def generate_ppt(self, HTML代码, 文件名, 默认字体="微软雅黑",
                     幻灯片宽度=12.8, 幻灯片高度=7.2, 下载地址前缀=""):
        prs = Presentation()
        prs.slide_width = PptInches(幻灯片宽度)
        prs.slide_height = PptInches(幻灯片高度)
        
        parser = SlideHTMLParser()
        parser.feed(HTML代码)
        slide_titles = [s.get("title", "") for s in parser.slides if s.get("title")]
        total_pages = len(parser.slides)

        for page_num, slide_data in enumerate(parser.slides, start=1):
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            theme = self._theme(slide_data.get("theme") or parser.global_theme)
            layout = (slide_data.get("layout") or "default").strip().lower()

            bg_color = self._resolve_color(slide_data.get("background"), theme["background"])
            if bg_color:
                self._draw_rect(slide, 0, 0, 幻灯片宽度, 幻灯片高度, bg=bg_color, border=bg_color, rounded=False)

            if layout in {"two-column", "comparison", "kpi-dashboard"}:
                lane_y = {"main": theme["margin_top"], "left": theme["margin_top"], "right": theme["margin_top"]}
            else:
                lane_y = {"main": theme["margin_top"], "left": theme["margin_top"], "right": theme["margin_top"]}
            section_state = {}

            if slide_data.get("title"):
                title_attrs = slide_data.get("title_style", {}) or {}
                title_style = self.parse_style(title_attrs.get("style", ""))
                title_font_size = self.get_font_size(title_style, 36)
                title_color = self.get_color(title_style, theme["title"])
                title_x = self._parse_length(self._style_or_attr(title_attrs, title_style, "x"), 0.5, 幻灯片宽度)
                title_y = self._parse_length(self._style_or_attr(title_attrs, title_style, "y"), theme["margin_top"], 幻灯片高度)
                title_w = self._parse_length(self._style_or_attr(title_attrs, title_style, "width"), 幻灯片宽度 - 1.0, 幻灯片宽度)
                title_h = self._parse_length(self._style_or_attr(title_attrs, title_style, "height"), 0.8, 幻灯片高度)
                self._add_textbox(
                    slide, title_x, title_y, title_w, title_h, slide_data["title"], 默认字体,
                    size=title_font_size, color=title_color, bold=True,
                    align=self._text_align(self._style_or_attr(title_attrs, title_style, "text-align"), PP_ALIGN.CENTER)
                )
                next_y = title_y + title_h + 0.2
                lane_y["main"] = max(lane_y["main"], next_y)
                lane_y["left"] = max(lane_y["left"], next_y)
                lane_y["right"] = max(lane_y["right"], next_y)
            
            for content in slide_data.get("content", []):
                content_type = content.get("type")
                lane = content.get("region") if content.get("region") in ("left", "right") else "main"
                default_x, default_w = self._region_box(lane, 幻灯片宽度)
                default_y = lane_y[lane]
                attrs = content.get("style", {}) or {}
                style = self.parse_style(attrs.get("style", ""))
                section = content.get("section")
                section_box = self._ensure_section_state(
                    section_state, section, slide, lane_y, lane, 幻灯片宽度, 幻灯片高度, theme
                ) if section else None

                if section_box:
                    lane = section_box["lane"]
                    default_x = section_box["x"] + section_box["padding"]
                    default_w = max(0.2, section_box["w"] - section_box["padding"] * 2)
                    default_y = section_box["cursor"]
                
                if content_type == "spacer":
                    sp_attrs = content.get("attrs", {}) or {}
                    sp_style = self.parse_style(sp_attrs.get("style", ""))
                    height = self._parse_length(self._style_or_attr(sp_attrs, sp_style, "height"), 0.22, 幻灯片高度)
                    if section_box:
                        section_box["cursor"] += height
                        lane_y[lane] = max(lane_y[lane], section_box["cursor"])
                    else:
                        lane_y[lane] += height
                    continue
                
                if content_type in ["h1", "h2"]:
                    default_sizes = {"h1": 28, "h2": 24}
                    font_size = self.get_font_size(style, default_sizes.get(content_type, 24))
                    color = self.get_color(style, theme["text"])
                    x = self._parse_length(self._style_or_attr(attrs, style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(attrs, style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(attrs, style, "width"), default_w, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(attrs, style, "height"), 0.7, 幻灯片高度)
                    self._add_textbox(
                        slide, x, y, w, h, content.get("text", ""), 默认字体, size=font_size,
                        color=color, bold=True, align=self._text_align(self._style_or_attr(attrs, style, "text-align"))
                    )
                    next_y = y + h + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y
                
                elif content_type == "p":
                    font_size = self.get_font_size(style, 18)
                    color = self.get_color(style, theme["text"])
                    x = self._parse_length(self._style_or_attr(attrs, style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(attrs, style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(attrs, style, "width"), default_w, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(attrs, style, "height"), 0.6, 幻灯片高度)
                    self._add_textbox(
                        slide, x, y, w, h, content.get("text", ""), 默认字体, size=font_size,
                        color=color, align=self._text_align(self._style_or_attr(attrs, style, "text-align"))
                    )
                    next_y = y + h + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y
                
                elif content_type == "list":
                    items = content.get("items", [])
                    ordered = content.get("ordered", False)
                    list_height_default = max(0.5, 0.4 * max(1, len(items)))
                    x = self._parse_length(self._style_or_attr(attrs, style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(attrs, style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(attrs, style, "width"), default_w, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(attrs, style, "height"), list_height_default, 幻灯片高度)
                    
                    text_box = slide.shapes.add_textbox(
                        PptInches(x), PptInches(y),
                        PptInches(w), PptInches(h)
                    )
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    
                    for i, item in enumerate(items):
                        item_attrs = item.get("style", {}) or {}
                        item_style = self.parse_style(item_attrs.get("style", ""))
                        font_size = self.get_font_size(item_style, 16)
                        color = self.get_color(item_style, theme["text"])
                        
                        if i == 0:
                            para = text_frame.paragraphs[0]
                        else:
                            para = text_frame.add_paragraph()
                        
                        prefix = f"{i+1}. " if ordered else "• "
                        para.text = prefix + item.get("text", "")
                        para.font.name = 默认字体
                        para.font.size = PptPt(font_size)
                        para.level = 0
                        para.alignment = self._text_align(self._style_or_attr(item_attrs, item_style, "text-align"))
                        if color:
                            color_tuple = self.hex_to_rgb(color)
                            r, g, b = color_tuple[:3]
                            para.font.color.rgb = RGBColor(r, g, b)
                            if OxmlElement is not None and len(color_tuple) > 3 and color_tuple[3] < 255:
                                opacity = int((color_tuple[3] / 255.0) * 100000)
                                srgbClr = para.font.color._xClr
                                alpha_elem = OxmlElement('a:alpha')
                                alpha_elem.set('val', str(opacity))
                                srgbClr.append(alpha_elem)
                    
                    next_y = y + h + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y

                elif content_type == "card":
                    card_attrs = content.get("attrs", {}) or {}
                    card_style = self.parse_style(card_attrs.get("style", ""))
                    x = self._parse_length(self._style_or_attr(card_attrs, card_style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(card_attrs, card_style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(card_attrs, card_style, "width"), default_w, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(card_attrs, card_style, "height"), 2.0, 幻灯片高度)
                    bg = self._resolve_color(
                        self._style_or_attr(card_attrs, card_style, "bg",
                                            self._style_or_attr(card_attrs, card_style, "background-color", theme["card_bg"])),
                        theme["card_bg"]
                    )
                    border = self._resolve_color(
                        self._style_or_attr(card_attrs, card_style, "border", theme["card_border"]),
                        theme["card_border"]
                    )
                    title_color = self._component_color(card_attrs, card_style, ["title-color", "titleColor"], theme["title"])
                    body_color = self._component_color(card_attrs, card_style, ["body-color", "text-color", "color", "bodyColor"], theme["text"])
                    title_text = card_attrs.get("title", "")
                    body_text = card_attrs.get("body", card_attrs.get("text", ""))
                    shadow = self._to_bool(self._style_or_attr(card_attrs, card_style, "shadow"), False)
                    rounded = self._to_bool(self._style_or_attr(card_attrs, card_style, "rounded", "true"), True)
                    radius = self._parse_length(self._style_or_attr(card_attrs, card_style, "radius"), None)
                    self._draw_rect(slide, x, y, w, h, bg=bg, border=border, rounded=rounded, shadow=shadow, radius=radius)

                    if title_text:
                        self._add_textbox(
                            slide, x + 0.2, y + 0.15, max(0.1, w - 0.4), 0.45, title_text,
                            默认字体, size=18, color=title_color, bold=True
                        )

                    if body_text:
                        self._add_textbox(
                            slide, x + 0.2, y + 0.65, max(0.1, w - 0.4), max(0.2, h - 0.8),
                            body_text.replace("&#10;", "\n"), 默认字体, size=14, color=body_color
                        )

                    next_y = y + h + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y

                elif content_type == "badge":
                    badge_attrs = content.get("attrs", {}) or {}
                    badge_style = self.parse_style(badge_attrs.get("style", ""))
                    x = self._parse_length(self._style_or_attr(badge_attrs, badge_style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(badge_attrs, badge_style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(badge_attrs, badge_style, "width"), 1.2, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(badge_attrs, badge_style, "height"), 0.45, 幻灯片高度)
                    bg = self._resolve_color(self._style_or_attr(badge_attrs, badge_style, "bg", "#E2E8F0"), "#E2E8F0")
                    color = self._resolve_color(self._style_or_attr(badge_attrs, badge_style, "color", theme["title"]), theme["title"])
                    text = badge_attrs.get("text", "BADGE")
                    font_size = int(self._parse_length(self._style_or_attr(badge_attrs, badge_style, "font-size"), 12))

                    shadow = self._to_bool(self._style_or_attr(badge_attrs, badge_style, "shadow"), False)
                    rounded = self._to_bool(self._style_or_attr(badge_attrs, badge_style, "rounded", "true"), True)
                    radius = self._parse_length(self._style_or_attr(badge_attrs, badge_style, "radius"), None)
                    shape = self._draw_rect(slide, x, y, w, h, bg=bg, border=bg, rounded=rounded, shadow=shadow, radius=radius)
                    tf = shape.text_frame
                    tf.clear()
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    para = tf.paragraphs[0]
                    para.text = text
                    para.alignment = PP_ALIGN.CENTER
                    para.font.name = 默认字体
                    para.font.bold = True
                    para.font.size = PptPt(font_size)
                    color_tuple = self.hex_to_rgb(color)
                    cr, cg, cb = color_tuple[:3]
                    para.font.color.rgb = RGBColor(cr, cg, cb)
                    if OxmlElement is not None and len(color_tuple) > 3 and color_tuple[3] < 255:
                        opacity = int((color_tuple[3] / 255.0) * 100000)
                        srgbClr = para.font.color._xClr
                        alpha_elem = OxmlElement('a:alpha')
                        alpha_elem.set('val', str(opacity))
                        srgbClr.append(alpha_elem)
                    next_y = y + h + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y

                elif content_type == "button":
                    btn_attrs = content.get("attrs", {}) or {}
                    btn_style = self.parse_style(btn_attrs.get("style", ""))
                    x = self._parse_length(self._style_or_attr(btn_attrs, btn_style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(btn_attrs, btn_style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(btn_attrs, btn_style, "width"), 2.0, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(btn_attrs, btn_style, "height"), 0.55, 幻灯片高度)
                    bg = self._resolve_color(self._style_or_attr(btn_attrs, btn_style, "bg", theme["primary"]), theme["primary"])
                    color = self._resolve_color(self._style_or_attr(btn_attrs, btn_style, "color", "#FFFFFF"), "#FFFFFF")
                    border = self._resolve_color(self._style_or_attr(btn_attrs, btn_style, "border", bg), bg)
                    text = btn_attrs.get("text", "按钮")
                    font_size = int(self._parse_length(self._style_or_attr(btn_attrs, btn_style, "font-size"), 14))

                    shadow = self._to_bool(self._style_or_attr(btn_attrs, btn_style, "shadow"), False)
                    rounded = self._to_bool(self._style_or_attr(btn_attrs, btn_style, "rounded", "true"), True)
                    radius = self._parse_length(self._style_or_attr(btn_attrs, btn_style, "radius"), None)
                    shape = self._draw_rect(slide, x, y, w, h, bg=bg, border=border, rounded=rounded, shadow=shadow, radius=radius)

                    tf = shape.text_frame
                    tf.clear()
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    para = tf.paragraphs[0]
                    para.text = text
                    para.alignment = PP_ALIGN.CENTER
                    para.font.name = 默认字体
                    para.font.bold = True
                    para.font.size = PptPt(font_size)
                    color_tuple = self.hex_to_rgb(color)
                    cr, cg, cb = color_tuple[:3]
                    para.font.color.rgb = RGBColor(cr, cg, cb)
                    if OxmlElement is not None and len(color_tuple) > 3 and color_tuple[3] < 255:
                        opacity = int((color_tuple[3] / 255.0) * 100000)
                        srgbClr = para.font.color._xClr
                        alpha_elem = OxmlElement('a:alpha')
                        alpha_elem.set('val', str(opacity))
                        srgbClr.append(alpha_elem)
                    next_y = y + h + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y

                elif content_type == "icon":
                    icon_attrs = content.get("attrs", {}) or {}
                    icon_style = self.parse_style(icon_attrs.get("style", ""))
                    x = self._parse_length(self._style_or_attr(icon_attrs, icon_style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(icon_attrs, icon_style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(icon_attrs, icon_style, "width"), default_w, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(icon_attrs, icon_style, "height"), 0.5, 幻灯片高度)
                    name = (icon_attrs.get("name") or "idea").strip().lower()
                    icon_char = self.ICON_MAP.get(name, "•")
                    text = icon_attrs.get("text", "")
                    color = self._resolve_color(self._style_or_attr(icon_attrs, icon_style, "color", theme["primary"]), theme["primary"])
                    body = f"{icon_char} {text}".strip()
                    self._add_textbox(slide, x, y, w, h, body, 默认字体, size=15, color=color)
                    next_y = y + h + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y

                elif content_type == "divider":
                    div_attrs = content.get("attrs", {}) or {}
                    div_style = self.parse_style(div_attrs.get("style", ""))
                    x = self._parse_length(self._style_or_attr(div_attrs, div_style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(div_attrs, div_style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(div_attrs, div_style, "width"), default_w, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(div_attrs, div_style, "height"), 0.03, 幻灯片高度)
                    color = self._resolve_color(self._style_or_attr(div_attrs, div_style, "color", theme["card_border"]), theme["card_border"])
                    self._draw_rect(slide, x, y, w, h, bg=color, border=color, rounded=False)
                    next_y = y + max(h, 0.03) + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y

                elif content_type == "callout":
                    co_attrs = content.get("attrs", {}) or {}
                    co_style = self.parse_style(co_attrs.get("style", ""))
                    x = self._parse_length(self._style_or_attr(co_attrs, co_style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(co_attrs, co_style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(co_attrs, co_style, "width"), default_w, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(co_attrs, co_style, "height"), 1.3, 幻灯片高度)
                    level = (co_attrs.get("level") or "info").lower()
                    palette = {
                        "info": ("#E0F2FE", "#0284C7"),
                        "success": ("#DCFCE7", "#16A34A"),
                        "warn": ("#FEF3C7", "#D97706"),
                        "error": ("#FEE2E2", "#DC2626")
                    }
                    bg, border = palette.get(level, palette["info"])
                    bg = self._resolve_color(self._style_or_attr(co_attrs, co_style, "bg", bg), bg)
                    border = self._resolve_color(self._style_or_attr(co_attrs, co_style, "border", border), border)
                    title_color = self._component_color(co_attrs, co_style, ["title-color", "titleColor", "header-color"], border)
                    body_color = self._component_color(co_attrs, co_style, ["body-color", "text-color", "color", "bodyColor"], theme["text"])
                    title = co_attrs.get("title", "")
                    body = co_attrs.get("body", co_attrs.get("text", ""))
                    shadow = self._to_bool(self._style_or_attr(co_attrs, co_style, "shadow"), False)
                    rounded = self._to_bool(self._style_or_attr(co_attrs, co_style, "rounded", "true"), True)
                    radius = self._parse_length(self._style_or_attr(co_attrs, co_style, "radius"), None)
                    self._draw_rect(slide, x, y, w, h, bg=bg, border=border, rounded=rounded, shadow=shadow, radius=radius)
                    if title:
                        self._add_textbox(slide, x + 0.18, y + 0.12, w - 0.3, 0.35, title, 默认字体, size=15, color=title_color, bold=True)
                    if body:
                        self._add_textbox(slide, x + 0.18, y + 0.5, w - 0.3, h - 0.58, body.replace("&#10;", "\n"), 默认字体, size=13, color=body_color)
                    next_y = y + h + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y

                elif content_type == "chart":
                    ch_attrs = content.get("attrs", {}) or {}
                    ch_style = self.parse_style(ch_attrs.get("style", ""))
                    x = self._parse_length(self._style_or_attr(ch_attrs, ch_style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(ch_attrs, ch_style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(ch_attrs, ch_style, "width"), default_w, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(ch_attrs, ch_style, "height"), 2.4, 幻灯片高度)
                    title = ch_attrs.get("title", "")
                    data = self._parse_pairs(ch_attrs.get("data", ""))
                    chart_type = (ch_attrs.get("type") or "bar").lower()
                    bar_color = self._resolve_color(ch_attrs.get("bar-color"), theme["accent"])
                    if title:
                        self._add_textbox(slide, x, y, w, 0.35, title, 默认字体, size=14, color=theme["title"], bold=True)
                    cy = y + (0.38 if title else 0.0)
                    ch = max(0.6, h - (0.38 if title else 0.0))
                    if chart_type == "progress":
                        value = 0.0
                        if data:
                            value = max(0.0, min(100.0, data[0][1]))
                        self._draw_rect(slide, x, cy + ch * 0.35, w, 0.22, bg="#E5E7EB", border="#E5E7EB", rounded=True)
                        self._draw_rect(slide, x, cy + ch * 0.35, w * (value / 100.0), 0.22, bg=bar_color, border=bar_color, rounded=True)
                        self._add_textbox(slide, x, cy + ch * 0.62, w, 0.3, f"{value:.1f}%", 默认字体, size=12, color=theme["muted"], align=PP_ALIGN.RIGHT)
                    else:
                        if not data:
                            data = [("A", 40), ("B", 60), ("C", 80)]
                        max_val = max(v for _, v in data) if data else 100.0
                        row_h = ch / max(1, len(data))
                        for i, (label, val) in enumerate(data):
                            row_y = cy + i * row_h
                            label_w = min(1.6, w * 0.3)
                            self._add_textbox(slide, x, row_y, label_w, row_h * 0.7, str(label), 默认字体, size=11, color=theme["text"])
                            bar_x = x + label_w + 0.08
                            bar_w = max(0.2, w - label_w - 0.2)
                            ratio = 0 if max_val <= 0 else max(0.0, min(1.0, val / max_val))
                            self._draw_rect(slide, bar_x, row_y + row_h * 0.2, bar_w, row_h * 0.4, bg="#E5E7EB", border="#E5E7EB", rounded=False)
                            self._draw_rect(slide, bar_x, row_y + row_h * 0.2, bar_w * ratio, row_h * 0.4, bg=bar_color, border=bar_color, rounded=False)
                            self._add_textbox(slide, bar_x + bar_w + 0.02, row_y, 0.5, row_h * 0.7, f"{val:g}", 默认字体, size=10, color=theme["muted"])
                    next_y = y + h + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y

                elif content_type == "table":
                    tb_attrs = content.get("attrs", {}) or {}
                    tb_style = self.parse_style(tb_attrs.get("style", ""))
                    x = self._parse_length(self._style_or_attr(tb_attrs, tb_style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(tb_attrs, tb_style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(tb_attrs, tb_style, "width"), default_w, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(tb_attrs, tb_style, "height"), 2.4, 幻灯片高度)
                    columns = [c.strip() for c in (tb_attrs.get("columns", "")).split(",") if c.strip()]
                    rows = self._parse_rows(tb_attrs.get("rows", ""))
                    col_count = max(1, len(columns), max((len(r) for r in rows), default=1))
                    row_count = 1 + len(rows)
                    cell_w = w / col_count
                    cell_h = h / row_count
                    header_bg = self._resolve_color(tb_attrs.get("header-bg"), theme["card_bg"])
                    stripe_bg = self._resolve_color(tb_attrs.get("stripe-bg"), "#F8FAFC")
                    border = self._resolve_color(tb_attrs.get("border"), theme["card_border"])
                    header_color = self._component_color(tb_attrs, tb_style, ["header-color", "title-color"], theme["title"])
                    text_color = self._component_color(tb_attrs, tb_style, ["text-color", "body-color", "color"], theme["text"])
                    for r in range(row_count):
                        for c in range(col_count):
                            cx = x + c * cell_w
                            cy = y + r * cell_h
                            bg = header_bg if r == 0 else (stripe_bg if r % 2 == 0 else "#FFFFFF")
                            self._draw_rect(slide, cx, cy, cell_w, cell_h, bg=bg, border=border, rounded=False)
                            if r == 0:
                                text = columns[c] if c < len(columns) else ""
                                tcolor = header_color
                                bold = True
                            else:
                                row = rows[r - 1] if r - 1 < len(rows) else []
                                text = row[c] if c < len(row) else ""
                                tcolor = text_color
                                bold = False
                            self._add_textbox(
                                slide, cx + 0.05, cy + 0.03, cell_w - 0.1, cell_h - 0.06,
                                text, 默认字体, size=11, color=tcolor, bold=bold, align=PP_ALIGN.CENTER
                            )
                    next_y = y + h + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y

                elif content_type == "toc":
                    toc_attrs = content.get("attrs", {}) or {}
                    toc_style = self.parse_style(toc_attrs.get("style", ""))
                    x = self._parse_length(self._style_or_attr(toc_attrs, toc_style, "x"), default_x, 幻灯片宽度)
                    y = self._parse_length(self._style_or_attr(toc_attrs, toc_style, "y"), default_y, 幻灯片高度)
                    w = self._parse_length(self._style_or_attr(toc_attrs, toc_style, "width"), default_w, 幻灯片宽度)
                    h = self._parse_length(self._style_or_attr(toc_attrs, toc_style, "height"), max(1.5, 幻灯片高度 - y - 0.8), 幻灯片高度)
                    titles = [t for t in slide_titles if t]
                    per_h = h / max(1, len(titles))
                    for i, t in enumerate(titles):
                        line_text = f"{i + 1}. {t}"
                        self._add_textbox(
                            slide, x, y + i * per_h, w, min(0.45, per_h),
                            line_text, 默认字体, size=16, color=theme["text"]
                        )
                    next_y = y + h + theme["line_gap"]
                    lane_y[lane] = max(lane_y[lane], next_y)
                    if section_box:
                        section_box["cursor"] = next_y

            self._render_footer(
                slide,
                slide_data.get("footer") or parser.global_footer,
                page_num,
                total_pages,
                幻灯片宽度,
                幻灯片高度,
                theme,
                默认字体,
                slide_data.get("title")
            )
        
        file_path = _get_unique_output_path(文件名, ".pptx")
        
        prs.save(file_path)
        download_url = _build_download_url(file_path, 下载地址前缀)
        return {
            "ui": {"text": [f"本地文件: {file_path}", f"下载链接: {download_url}"]},
            "result": (file_path, download_url)
        }
    
NODE_CLASS_MAPPINGS = {
    "TextToWord": TextToWordNode,
    "TextToWordAdvanced": TextToWordAdvancedNode,
    "HTMLToPPT": HTMLToPPTNode,
}

NODE_DISPLAY_NAME_MAPPINGS = {
    "TextToWord": "文字生成Word",
    "TextToWordAdvanced": "文字生成Word(高级)",
    "HTMLToPPT": "HTML生成PPT",
}
